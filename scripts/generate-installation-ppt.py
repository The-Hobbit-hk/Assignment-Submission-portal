"""
Generate the master Opulence DRR & Council Installation PPT.

Royal maroon + gold "Opulence" theme to match the event PR creatives.

  python scripts/generate-installation-ppt.py

Output: output/Opulence-DRR-Council-Installation-2026-27.pptx
"""
from __future__ import annotations

import re
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parent.parent
PUBLIC = ROOT / "public"
OUT_DIR = ROOT / "output"
OUT_FILE = OUT_DIR / "Opulence-DRR-Council-Installation-2026-27.pptx"

# ── Opulence palette (maroon + gold) ─────────────────────────────────────────
MAROON_DEEP = RGBColor(0x33, 0x05, 0x0B)   # near-black maroon (vignette)
MAROON = RGBColor(0x5C, 0x0A, 0x14)         # primary royal maroon
MAROON_LT = RGBColor(0x7C, 0x14, 0x20)      # lighter maroon for gradient
GOLD = RGBColor(0xC9, 0xA2, 0x4B)           # antique gold
GOLD_BRIGHT = RGBColor(0xE8, 0xC8, 0x6A)    # bright gold highlight
GOLD_PALE = RGBColor(0xEA, 0xD0, 0x8A)      # pale gold
CREAM = RGBColor(0xF5, 0xEC, 0xD6)          # parchment / plaque fill
CREAM_DK = RGBColor(0xE9, 0xDC, 0xBE)       # alt row parchment
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x2A, 0x0A, 0x0E)            # deep maroon ink for text on cream

TITLE_FONT = "Book Antiqua"                  # elegant serif for headings
BODY_FONT = "Georgia"                        # serif body

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Council roster (matches council-roster-data.ts) ──────────────────────────
COUNCIL = [
    ("PHF. DRR. Dr. Karishma Awari", "District Rotaract Representative", "Rotaract Club of Pune Shaniwarwada", "karishma-awari.png"),
    ("PHF. PDRR. Drishti Singh", "District Learning Facilitator", "Rotaract Club of Humanitas", None),
    ("PHF. Rtr. Harshvardhan Kale", "District General Secretary", "Rotaract Club of Bavdhan Pioneers", "harshvardhan-kale.png"),
    ("Rtr. Suraj Surkutla", "District Secretary - Administration", "Rotaract Club of Pune Baner", "suraj-surkutla.png"),
    ("Rtr. Hamid Shaikh", "District Secretary - Events", "Rotaract Club of Aundh Smartcity", "hamid-shaikh.png"),
    ("Rtr. Disha Daga", "District Secretary - Protocols", "Rotaract Club of Bibwewadi Pune", "disha-daga.png"),
    ("Rtr. Dr. Aishwarya Patil", "District Secretary - Reporting", "Rotaract Club of Pune Shaniwarwada", "aishwarya-patil.png"),
    ("PHF. Rtr. Sharvindu Jogdand", "District Treasurer", "Rotaract Club of Pune Warje", "sharvindu-jogdand.png"),
    ("Rtr. Dr. Ashlesha Deshpande", "District Club Advisor", "Rotaract Club of Pune Heritage", "ashlesha-deshpande.png"),
    ("Rtr. Rohan Puri", "Zonal Advisor", "Rotaract Club of Khopoli", "rohan-puri.png"),
    ("Rtr. Aniket Sardar", "District Zonal Representative", "Rotaract Club of Khopoli", "aniket-sardar.png"),
    ("Rtr. Vedant Chirmade", "District Zonal Representative", "Rotaract Club of Pimpri", "vedant-chirmade.png"),
    ("Rtr. Vedant Chaudhari", "District Zonal Representative", "Rotaract Club of Pimpri", "vedant-chaudhari.png"),
    ("Rtr. Tisha Sancheti", "District Zonal Representative", "Rotaract Club of Pune Camp Next Gen", "tisha-sancheti.png"),
    ("Rtr. Pratham Pokharkar", "District Zonal Representative", "Rotaract Club of Pune Aurora", "pratham-pokharkar.png"),
    ("Rtr. Rajas Uchagaonkar", "District Zonal Representative", "Rotaract Club of Pune City Fortune", "rajas-uchagaonkar.png"),
    ("Rtr. Prem Bansode", "District Zonal Representative", "Rotaract Club of Daund College", "prem-bansode.png"),
    ("Rtr. Prerna Bhilare", "Assistant Zonal Representative", "Rotaract Club of Sinhgad College Of Pharmacy", "prerna-bhilare.png"),
    ("Rtr. Aditya Verma", "Assistant Zonal Representative", "Rotaract Club of Nigdi Pune", "aditya-verma.png"),
    ("Rtr. Shrushti Shirore", "Assistant Zonal Representative", "Rotaract Club of Pune Aurora", "shrushti-shirore.png"),
    ("Rtr. Sarthak Ambhore", "Assistant Zonal Representative", "Rotaract Club of Pimpri", "sarthak-ambhore.png"),
    ("Rtr. Rohit Kumbhar", "Assistant Zonal Representative", "Rotaract Club of Bavdhan Pioneers", "rohit-kumbhar.png"),
    ("Rtr. Sumedh Gite", "Assistant Zonal Representative", "Rotaract Club of Aundh Smartcity", "sumedh-gite.png"),
    ("Rtr. Samrudhi Khade", "District Director - Professional Development", "Rotaract Club of Pune Zenith", "samrudhi-khade.png"),
    ("Rtr. Jayesh Chavan", "District Director - Club Service", "Rotaract Club of Pune City Fortune", "jayesh-chavan.png"),
    ("PHF. Rtr. Aslam Dhanani", "District Director - Community Service", "Rotaract Club of Aundh", "aslam-dhanani.png"),
    ("Rtr. Vaishnavi Kedari", "District Co-Director - Community Service", "Rotaract Club of Symbiosis Skills and Professional University", "vaishnavi-kedari.png"),
    ("PHF. Rtr. Ishan Malawade", "District Director - International Service", "Rotaract Club of Vibrants", "ishan-malawade.png"),
    ("Rtr. Pranav Gandhi", "District Co-Director - International Service", "Rotaract Club of Bibwewadi Pune", "pranav-gandhi.png"),
    ("Rtr. Omkar Patil", "District Director - Membership Development", "Rotaract Club of Pune City Fortune", "omkar-patil.png"),
    ("Rtr. Chinmayee Bartakke", "District Director - Diversity, Equity & Inclusion", "Rotaract Club of Viman Nagar", "chinmayee-bartakke.png"),
    ("Rtr. Faizan Tamboli", "District Director - Communications", "Rotaract Club of Pune Shaniwarwada", "faizan-tamboli.png"),
    ("Rtr. Shrawani Shendkar", "District Director - Public Image", "Rotaract Club of Genba Sopanrao Moze College of Engineering", "shrawani-shendkar.png"),
    ("Rtr. Salvin Padvi", "District Officer - Chief Branding", "Rotaract Club of Pune Aurora", "salvin-padvi.png"),
    ("Rtr. Janhavi Yeole", "District Officer - Public Relations", "Rotaract Club of Pune Zenith", "janhavi-yeole.png"),
    ("Rtr. Shreeraj Nilkanth", "District Officer - Public Relations", "Rotaract Club of Panvel Industrial Town", "shreeraj-nilkanth.png"),
    ("Rtr. Sushant Chavan", "District Officer - Editing", "Rotaract Club of Sinhgad College Of Pharmacy", "sushant-chavan.png"),
    ("Rtr. Harshal Nikam", "District Officer - Editing", "Rotaract Club of Pune Heritage", "harshal-nikam.png"),
    ("Rtr. Abhishek Dixit", "District Officer - Editing", "Rotaract Club of Vibrants", "abhishek-dixit.png"),
    ("Rtr. Vageesha Karhadkar", "District Officer - Editing", "Rotaract Club of Magarpatta Trendsetters", "vageesha-karhadkar.png"),
    ("Rtr. Shreyas Pathak", "District Sergeant-at-Arms", "Rotaract Club of Pune Mideast", "shreyas-pathak.png"),
    ("Rtr. Snehal Jadhav", "District Sergeant-at-Arms", "Rotaract Club of Balewadi High Street", "snehal-jadhav.png"),
    ("Rtr. Pranav Pisal", "District Officer - Grants & RYLA", "Rotaract Club of Pimpri", "pranav-pisal.png"),
    ("Rtr. Adhishree Thakar", "District Officer - Rotary Rotaract Relations", "Rotaract Club of Pune Zenith", "adhishree-thakar.png"),
    ("Rtr. Prajwal Bande", "District Officer - Interact Rotaract Relations", "Rotaract Club of Daund College", "prajwal-bande.png"),
    ("Rtr. Gaurav Golecha", "District Officer - Professional Assistance", "Rotaract Club of Pune Mideast", "gaurav-golecha.png"),
    ("Rtr. Talha Shaikh", "District Officer - Professional Assistance", "Rotaract Club of Aundh Smartcity", "talha-shaikh.png"),
    ("PHF. DRRE. Adv. Sattyajeet Karale Patil", "District Legal Advisor", "Rotaract Club of Pune Samrajya", "sattyajeet-karale-patil.png"),
    ("Rtr. Devsharan Singh", "District Coordinator - Website", "Rotaract Club of Aundh Smartcity", "devsharan-singh.png"),
    ("Rtr. Ashi Agarwal", "District Coordinator - Multi District Events", "Rotaract Club of Roar NIBM", "ashi-agarwal.png"),
    ("Rtr. Sanjana Pawar", "District Chairperson - World Rotaract Week", "Rotaract Club of Vibrants", "sanjana-pawar.png"),
    ("Rtr. Vedant Buge", "District Officer - Without Portfolio", "Rotaract Club of Pune Kalyani Nagar", "vedant-buge.png"),
    ("Rtr. Priya Bhagwani", "District Officer - Without Portfolio", "Rotaract Club of Nigdi Pune", "priya-bhagwani.png"),
    ("PHF. Rtr. Vansh Chawla", "District Convenor - District Sports Meet", "Rotaract Club of Pimpri", "vansh-chawla.png"),
    ("Rtr. Amruta Potdukhe", "District Convenor - DRR and Council Installation", "Rotaract Club of Sinhgad College of Pharmacy", "amruta-potdukhe.png"),
    ("Rtr. Digvijay Lad", "District Convenor - District Trek", "Rotaract Club of Pune City Fortune", "digvijay-lad.png"),
    ("Rtr. Vijeta Kulkarni", "District Convenor - District Culturals", "Rotaract Club of Pune Samrajya", "vijeta-kulkarni.png"),
]

DGS_NAME = "PHF. Rtr. Harshvardhan Kale"
LOGO = PUBLIC / "logo-rotaract-3131.png"
LOGO_MARK = PUBLIC / "logo-rotaract-mark.png"
REIGN = PUBLIC / "reign-theme-riy-2026-27.png"


# ── Low-level helpers ────────────────────────────────────────────────────────
def photo_path(filename: str | None) -> Path | None:
    if not filename:
        return None
    p = PUBLIC / "council" / filename
    return p if p.exists() else None


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _resolve_out(path: Path) -> Path:
    """Fall back to a versioned name if the target file is locked/open."""
    if not path.exists():
        return path
    try:
        with open(path, "a"):
            return path
    except PermissionError:
        i = 2
        while True:
            alt = path.with_name(f"{path.stem}-v{i}{path.suffix}")
            if not alt.exists():
                return alt
            i += 1


def solid_bg(slide, color: RGBColor) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _apply_gradient(shape, c1: RGBColor, c2: RGBColor, angle: int = 90) -> None:
    shape.fill.gradient()
    stops = shape.fill.gradient_stops
    stops[0].position = 0.0
    stops[0].color.rgb = c1
    stops[1].position = 1.0
    stops[1].color.rgb = c2
    try:
        shape.fill.gradient_angle = angle
    except Exception:
        pass


def _send_to_back(slide, shape) -> None:
    spTree = slide.shapes._spTree
    spTree.remove(shape._element)
    spTree.insert(2, shape._element)


def add_shape(slide, shape_type, left, top, width, height, color=None, line=None, line_w=1.5):
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    if color is not None:
        shp.fill.solid()
        shp.fill.fore_color.rgb = color
    else:
        shp.fill.background()
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def add_text(slide, left, top, width, height, text, size=24, bold=False,
             color=INK, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE,
             spacing=None, font=BODY_FONT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        run.font.name = font
        p.alignment = align
        if spacing is not None:
            run.font._rPr.set("spc", str(int(spacing * 100)))
    return box


# ── Themed decoration ────────────────────────────────────────────────────────
def maroon_bg(slide) -> None:
    """Royal maroon vignette gradient behind everything."""
    shp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    shp.line.fill.background()
    _apply_gradient(shp, MAROON_LT, MAROON_DEEP, angle=90)
    _send_to_back(slide, shp)


def gold_frame(slide, inset=Inches(0.28)) -> None:
    """Double ornamental gold border + corner diamonds."""
    outer = add_shape(slide, MSO_SHAPE.RECTANGLE, inset, inset,
                      Emu(SLIDE_W - 2 * inset), Emu(SLIDE_H - 2 * inset),
                      color=None, line=GOLD, line_w=2.5)
    inner_in = Emu(inset + Inches(0.09))
    add_shape(slide, MSO_SHAPE.RECTANGLE, inner_in, inner_in,
              Emu(SLIDE_W - 2 * inner_in), Emu(SLIDE_H - 2 * inner_in),
              color=None, line=GOLD, line_w=0.75)
    # corner diamonds
    d = Inches(0.16)
    for cx, cy in ((inset, inset),
                   (Emu(SLIDE_W - inset), inset),
                   (inset, Emu(SLIDE_H - inset)),
                   (Emu(SLIDE_W - inset), Emu(SLIDE_H - inset))):
        add_shape(slide, MSO_SHAPE.DIAMOND, Emu(cx - d / 2), Emu(cy - d / 2), d, d, GOLD_BRIGHT)


def gold_rule(slide, cx, y, half_w=Inches(1.6)) -> None:
    """Gold divider line with a centred diamond."""
    add_shape(slide, MSO_SHAPE.RECTANGLE, Emu(cx - half_w), y, Emu(half_w - Inches(0.22)), Pt(1.6), GOLD)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Emu(cx + Inches(0.22)), y, Emu(half_w - Inches(0.22)), Pt(1.6), GOLD)
    d = Inches(0.16)
    add_shape(slide, MSO_SHAPE.DIAMOND, Emu(cx - d / 2), Emu(y - d / 2 + Pt(0.8)), d, d, GOLD_BRIGHT)


def cream_plaque(slide, left, top, width, height, text, size=15, bold=True,
                 spacing=1.0, align=PP_ALIGN.CENTER):
    """Parchment label box with gold border (matches PR sub-title plaque)."""
    plq = add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height, CREAM,
                    line=GOLD, line_w=1.25)
    add_text(slide, left, top, width, height, text, size=size, bold=bold,
             color=MAROON, align=align, spacing=spacing, font=TITLE_FONT)
    return plq


def brand_footer(slide, time_slot: str) -> None:
    add_shape(slide, MSO_SHAPE.RECTANGLE, Inches(0.75), Inches(6.78), Inches(11.83), Pt(1), GOLD)
    add_text(slide, Inches(0.75), Inches(6.85), Inches(8), Inches(0.35),
             "OPULENCE  ·  DRR & Council Installation  ·  RIY 2026-27",
             size=9, color=GOLD_PALE, align=PP_ALIGN.LEFT, font=TITLE_FONT, spacing=0.5)
    if time_slot:
        add_text(slide, Inches(5.0), Inches(6.85), Inches(7.58), Inches(0.35),
                 time_slot, size=9, bold=True, color=GOLD_BRIGHT, align=PP_ALIGN.RIGHT, font=TITLE_FONT)


def brand_header(slide) -> None:
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Inches(0.7), Inches(0.5), height=Inches(0.5))
    add_text(slide, Inches(1.35), Inches(0.52), Inches(6), Inches(0.45),
             "ROTARACT DISTRICT 3131", size=11, bold=True, color=GOLD_PALE,
             align=PP_ALIGN.LEFT, spacing=2.0, font=TITLE_FONT)


def _add_circular_picture(slide, img_path: str, left, top, size) -> None:
    pic = slide.shapes.add_picture(img_path, left, top, size, size)
    spPr = pic._element.spPr
    for tag in ("a:prstGeom", "a:custGeom"):
        ex = spPr.find(qn(tag))
        if ex is not None:
            spPr.remove(ex)
    geom = spPr.makeelement(qn("a:prstGeom"), {"prst": "ellipse"})
    geom.append(spPr.makeelement(qn("a:avLst"), {}))
    spPr.append(geom)


def _initials(name: str) -> str:
    cleaned = re.sub(r"^(PHF\.|Rtr\.|Dr\.|DRR\.|PDRR\.|DRRE\.|Adv\.|Rtn\.|\s)+", "", name).strip()
    parts = [p for p in cleaned.split() if p]
    return "".join(w[0] for w in parts[:2]).upper() or "R"


# ── Slide builders ───────────────────────────────────────────────────────────
def slide_cover(prs: Presentation) -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(5.17), Inches(0.72), width=Inches(3.0))
    add_text(slide, Inches(0.5), Inches(1.95), Inches(12.33), Inches(1.5), "OPULENCE",
             size=88, bold=True, color=GOLD_BRIGHT, spacing=6.0, font=TITLE_FONT)
    cream_plaque(slide, Inches(3.66), Inches(3.55), Inches(6.0), Inches(0.95),
                 "District Assembly\nDRR & Council Installation", size=17, spacing=0.5)
    gold_rule(slide, Emu(SLIDE_W / 2), Inches(4.85))
    add_text(slide, Inches(0.5), Inches(5.05), Inches(12.33), Inches(0.5),
             "Rotaract District 3131  ·  Rotary International Year 2026-27",
             size=16, color=CREAM, font=TITLE_FONT, italic=True)
    add_text(slide, Inches(0.5), Inches(5.75), Inches(12.33), Inches(0.5),
             "4th July 2026  ·  4:00 PM onwards", size=15, bold=True, color=GOLD_PALE,
             font=TITLE_FONT, spacing=1.0)
    add_text(slide, Inches(0.5), Inches(6.2), Inches(12.33), Inches(0.5),
             "SAAFA Banquets, Balewadi High Street, Baner", size=13, color=CREAM, font=TITLE_FONT)


def slide_section(prs: Presentation, title: str, time_slot: str = "") -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    if LOGO_MARK.exists():
        slide.shapes.add_picture(str(LOGO_MARK), Emu(SLIDE_W / 2 - Inches(0.45)), Inches(1.15),
                                 height=Inches(0.9))
    n_lines = title.count("\n") + 1
    top = Inches(3.0) if n_lines > 1 else Inches(3.2)
    add_text(slide, Inches(0.8), top, Inches(11.73), Inches(1.6), title,
             size=44, bold=True, color=GOLD_BRIGHT, font=TITLE_FONT, spacing=1.0)
    gold_rule(slide, Emu(SLIDE_W / 2), Inches(5.15))
    if time_slot:
        add_text(slide, Inches(0.8), Inches(5.35), Inches(11.73), Inches(0.5), time_slot,
                 size=16, bold=True, color=CREAM, font=TITLE_FONT, spacing=1.5)
    brand_footer(slide, "")


def slide_moment(prs: Presentation, title: str, subtitle: str, time_slot: str, label: str = "Ceremony") -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    brand_header(slide)
    cream_plaque(slide, Emu(SLIDE_W / 2 - Inches(1.6)), Inches(2.3), Inches(3.2), Inches(0.5),
                 label.upper(), size=13, spacing=2.0)
    add_text(slide, Inches(0.8), Inches(3.05), Inches(11.73), Inches(1.3), title,
             size=46, bold=True, color=GOLD_BRIGHT, font=TITLE_FONT, spacing=0.5)
    gold_rule(slide, Emu(SLIDE_W / 2), Inches(4.5))
    if subtitle:
        add_text(slide, Inches(0.8), Inches(4.7), Inches(11.73), Inches(0.7), subtitle,
                 size=19, color=CREAM, font=TITLE_FONT, italic=True)
    brand_footer(slide, time_slot)


def slide_speaker(prs: Presentation, heading: str, name: str, role: str, time_slot: str,
                  photo_file: str | None = None) -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    brand_header(slide)

    path = photo_path(photo_file)
    text_left = Inches(1.05)
    text_w = Inches(11.2)
    if path:
        # gold double ring around circular photo
        ring = add_shape(slide, MSO_SHAPE.OVAL, Inches(8.55), Inches(1.95), Inches(3.5), Inches(3.5),
                         color=None, line=GOLD, line_w=3)
        add_shape(slide, MSO_SHAPE.OVAL, Inches(8.7), Inches(2.1), Inches(3.2), Inches(3.2),
                  color=None, line=GOLD_PALE, line_w=1)
        _add_circular_picture(slide, str(path), Inches(8.78), Inches(2.18), Inches(3.04))
        text_w = Inches(7.3)

    cream_plaque(slide, text_left, Inches(2.45),
                 min(Inches(5.4), Inches(0.15 * len(heading) + 1.0)), Inches(0.5),
                 heading.upper(), size=13, spacing=1.5, align=PP_ALIGN.CENTER)
    add_text(slide, text_left, Inches(3.15), text_w, Inches(1.5), name,
             size=38, bold=True, color=GOLD_BRIGHT, align=PP_ALIGN.LEFT, font=TITLE_FONT)
    if role:
        add_text(slide, Emu(text_left + Inches(0.03)), Inches(4.55), text_w, Inches(0.8), role,
                 size=19, bold=True, color=CREAM, align=PP_ALIGN.LEFT, font=TITLE_FONT, italic=True)
    brand_footer(slide, time_slot)


def slide_member(prs: Presentation, name: str, title: str, club: str, photo_file: str | None,
                 time_slot: str, index: int) -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    brand_header(slide)

    # faint gold index watermark
    add_text(slide, Inches(9.2), Inches(0.35), Inches(3.4), Inches(1.4), f"{index:02d}",
             size=90, bold=True, color=MAROON_LT, align=PP_ALIGN.RIGHT, font=TITLE_FONT)

    # Photo in ornate gold-framed cream card (left)
    card_l, card_t, card_w, card_h = Inches(0.95), Inches(1.55), Inches(3.9), Inches(4.5)
    add_shape(slide, MSO_SHAPE.RECTANGLE, Emu(card_l - Inches(0.12)), Emu(card_t - Inches(0.12)),
              Emu(card_w + Inches(0.24)), Emu(card_h + Inches(0.24)), color=None, line=GOLD, line_w=2.5)
    card = add_shape(slide, MSO_SHAPE.RECTANGLE, card_l, card_t, card_w, card_h, CREAM,
                     line=GOLD_PALE, line_w=0.75)
    path = photo_path(photo_file)
    if path:
        slide.shapes.add_picture(str(path), Emu(card_l + Inches(0.1)), Emu(card_t + Inches(0.1)),
                                 Emu(card_w - Inches(0.2)), Emu(card_h - Inches(0.2)))
    else:
        add_shape(slide, MSO_SHAPE.OVAL, Emu(card_l + Inches(1.05)), Emu(card_t + Inches(1.2)),
                  Inches(1.8), Inches(1.8), color=None, line=GOLD, line_w=2)
        add_text(slide, card_l, Emu(card_t + Inches(1.45)), card_w, Inches(1.3), _initials(name),
                 size=54, bold=True, color=MAROON, font=TITLE_FONT)
        add_text(slide, card_l, Emu(card_t + Inches(3.0)), card_w, Inches(0.6), "Photo awaited",
                 size=12, color=INK, font=TITLE_FONT, italic=True)

    # Right side content
    rx = Inches(5.5)
    rw = Inches(7.0)
    cream_plaque(slide, rx, Inches(1.75), Inches(3.5), Inches(0.5),
                 "OFFICIAL INSTALLATION", size=12, spacing=1.5)
    add_text(slide, rx, Inches(2.55), rw, Inches(1.5), name,
             size=32, bold=True, color=GOLD_BRIGHT, align=PP_ALIGN.LEFT, font=TITLE_FONT)
    # role band (gold fill, maroon text)
    role_w = min(Inches(6.9), Inches(0.155 * len(title) + 0.7))
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, rx, Inches(4.15), role_w, Inches(0.62), GOLD,
              line=GOLD_BRIGHT, line_w=1)
    add_text(slide, rx, Inches(4.15), role_w, Inches(0.62), title, size=15, bold=True,
             color=MAROON, align=PP_ALIGN.CENTER, font=TITLE_FONT)
    add_text(slide, Emu(rx + Inches(0.02)), Inches(5.05), rw, Inches(0.9), club,
             size=15, color=CREAM, align=PP_ALIGN.LEFT, font=TITLE_FONT, italic=True)
    brand_footer(slide, time_slot)


def slide_agenda(prs: Presentation, rows, page_title: str, part: str) -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    brand_header(slide)
    add_text(slide, Inches(0.85), Inches(1.15), Inches(9), Inches(0.7), page_title,
             size=30, bold=True, color=GOLD_BRIGHT, align=PP_ALIGN.LEFT, font=TITLE_FONT, spacing=0.5)
    add_text(slide, Inches(9.3), Inches(1.28), Inches(3.2), Inches(0.5), part,
             size=13, bold=True, color=GOLD_PALE, align=PP_ALIGN.RIGHT, font=TITLE_FONT, spacing=1.0)

    left0 = Inches(0.85)
    col_w = [Inches(7.55), Inches(1.6), Inches(2.6)]
    top = Inches(1.95)
    header_h = Inches(0.42)
    # header
    x = left0
    for i, h in enumerate(("Description", "Duration", "Time")):
        add_shape(slide, MSO_SHAPE.RECTANGLE, x, top, col_w[i], header_h, GOLD)
        add_text(slide, x, top, col_w[i], header_h, h, size=11, bold=True, color=MAROON, font=TITLE_FONT)
        x = Emu(x + col_w[i])
    y = Emu(top + header_h)
    row_h = Inches(0.38)
    for r, (desc, dur, time) in enumerate(rows):
        bg = CREAM if r % 2 == 0 else CREAM_DK
        x = left0
        for i, val in enumerate((desc, dur, time)):
            cell = add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, col_w[i], row_h, bg,
                             line=GOLD_PALE, line_w=0.5)
            add_text(slide, Emu(x + Inches(0.12)), y, Emu(col_w[i] - Inches(0.2)), row_h, val,
                     size=9.5 if i == 0 else 10, color=INK,
                     align=PP_ALIGN.LEFT if i == 0 else PP_ALIGN.CENTER, font=BODY_FONT)
            x = Emu(x + col_w[i])
        y = Emu(y + row_h)
    brand_footer(slide, "")


def slide_reign(prs: Presentation) -> None:
    slide = blank(prs)
    maroon_bg(slide)
    gold_frame(slide)
    brand_header(slide)
    if REIGN.exists():
        slide.shapes.add_picture(str(REIGN), Inches(4.9), Inches(1.25), height=Inches(1.8))
    add_text(slide, Inches(0.85), Inches(3.4), Inches(11.63), Inches(1.0),
             "Rotaract Empowering Individuals\nfor Growth and Networking",
             size=26, bold=True, color=GOLD_BRIGHT, font=TITLE_FONT)
    gold_rule(slide, Emu(SLIDE_W / 2), Inches(4.85))
    add_text(slide, Inches(0.85), Inches(5.05), Inches(11.63), Inches(0.6),
             "The REIGN Vision  ·  DRR Dr. Karishma Awari  ·  RIY 2026-27",
             size=15, color=CREAM, font=TITLE_FONT, italic=True)
    brand_footer(slide, "7:15 – 8:00 PM")


def find_member(fragment: str):
    for m in COUNCIL:
        if fragment.lower() in m[0].lower():
            return m
    return (fragment, "", "", None)


def council_except(*skip: str):
    s = {x.lower() for x in skip}
    return [m for m in COUNCIL if m[0].lower() not in s]


def main() -> None:
    OUT_DIR.mkdir(parents=True, exist_ok=True)
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_cover(prs)

    agenda_1 = [
        ("Registrations Open", "30 Mins", "4:30 – 5:00"),
        ("Commencement of Opulence: DRR & Council Installation", "—", "5:00 – 9:13 PM"),
        ("Ganesh Vandana", "5 mins", "5:00 – 5:05 PM"),
        ("Go-Green Initiative followed by Rotaract Protocols", "5 mins", "5:05 – 5:10 PM"),
        ("Welcome Address by Host Club Presidents", "2 mins", "5:10 – 5:12 PM"),
        ("Event Flow by DPS Rtr. Disha Daga", "2 mins", "5:12 – 5:14 PM"),
        ("Introduction of DRR Karishma Awari", "5 mins", "5:14 – 5:19 PM"),
        ("Official Installation of DRR Karishma Awari", "5 mins", "5:19 – 5:24 PM"),
        ("Official Address by DRR Karishma Awari", "15 mins", "5:24 – 5:39 PM"),
        ("Introduction of DGS Harshvardhan Kale", "2 mins", "5:39 – 5:41 PM"),
        ("Official Installation of DGS Harshvardhan Kale", "2 mins", "5:41 – 5:43 PM"),
        ("Official Address by DGS Harshvardhan Kale", "5 mins", "5:43 – 5:48 PM"),
    ]
    agenda_2 = [
        ("Official Installation of all District Council Members", "60 mins", "5:48 – 6:48 PM"),
        ("Oath Ceremony of Council", "2 mins", "6:48 – 6:50 PM"),
        ("Introduction of Esteemed Guest – DG Rtn. Nitin Dhamale", "3 mins", "6:50 – 6:53 PM"),
        ("Address by Esteemed Guest – DG Nitin Dhamale", "15 mins", "6:53 – 7:08 PM"),
        ("Felicitation of Esteemed Guest – DG Nitin Dhamale", "2 mins", "7:08 – 7:10 PM"),
        ("Address & Felicitation of Guest – DRCC Rtn. Balvir Chawla", "5 mins", "7:10 – 7:15 PM"),
        ("REIGN Vision", "45 mins", "7:15 – 8:00 PM"),
        ("Launch of District Website", "2 mins", "8:00 – 8:02 PM"),
        ("Launch of District Directory", "5 mins", "8:02 – 8:07 PM"),
        ("Special Announcements", "5 mins", "8:07 – 8:12 PM"),
        ("Sponsor Session", "7 mins", "8:12 – 8:19 PM"),
        ("Address & Felicitation of Guest – Dist. YSD Rtn. Santosh Pardeshi", "5 mins", "8:19 – 8:24 PM"),
    ]
    agenda_3 = [
        ("Address & Felicitation – Galaxy of DRRs", "15 min", "8:24 – 8:39 PM"),
        ("Introduction of Chief Guest – Past RID Dr. Mahesh Kotbagi", "2 mins", "8:39 – 8:41 PM"),
        ("Address by Chief Guest – Past RID Dr. Mahesh Kotbagi", "20 mins", "8:41 – 9:01 PM"),
        ("Felicitation of Chief Guest – Past RID Dr. Mahesh Kotbagi", "2 mins", "9:01 – 9:03 PM"),
        ("Closing Address by HC", "2 mins", "9:03 – 9:05 PM"),
        ("Secretarial Announcements", "2 mins", "9:05 – 9:07 PM"),
        ("Open Forum", "2 mins", "9:07 – 9:09 PM"),
        ("Vote of Thanks", "2 mins", "9:09 – 9:11 PM"),
        ("Opulence DRR & Council Installation adjourned", "2 mins", "9:11 – 9:13 PM"),
        ("Dinner", "—", "9:13 PM onwards"),
    ]
    slide_agenda(prs, agenda_1, "Event Agenda", "Part I of III")
    slide_agenda(prs, agenda_2, "Event Agenda", "Part II of III")
    slide_agenda(prs, agenda_3, "Event Agenda", "Part III of III")

    slide_moment(prs, "Registrations Open", "Welcome to Opulence", "4:30 – 5:00 PM", "Welcome")
    slide_section(prs, "Commencement of\nOpulence", "5:00 PM")

    slide_moment(prs, "Ganesh Vandana", "Invoking blessings for the evening", "5:00 – 5:05 PM", "Opening")
    slide_moment(prs, "Go-Green Initiative", "followed by Rotaract Protocols", "5:05 – 5:10 PM", "Protocols")
    slide_speaker(prs, "Welcome Address", "Host Club Presidents", "Opulence Host Committee", "5:10 – 5:12 PM")
    disha = find_member("Disha Daga")
    slide_speaker(prs, "Event Flow", disha[0], disha[1], "5:12 – 5:14 PM", disha[3])

    drr = find_member("Karishma Awari")
    slide_speaker(prs, "Introduction of", drr[0], drr[1], "5:14 – 5:19 PM", drr[3])
    slide_member(prs, drr[0], drr[1], drr[2], drr[3], "5:19 – 5:24 PM", 1)
    slide_speaker(prs, "Official Address", drr[0], drr[1], "5:24 – 5:39 PM", drr[3])

    dgs = find_member("Harshvardhan Kale")
    slide_speaker(prs, "Introduction of", dgs[0], dgs[1], "5:39 – 5:41 PM", dgs[3])
    slide_member(prs, dgs[0], dgs[1], dgs[2], dgs[3], "5:41 – 5:43 PM", 2)
    slide_speaker(prs, "Official Address", dgs[0], dgs[1], "5:43 – 5:48 PM", dgs[3])

    slide_section(prs, "Official Installation of\nAll District Council Members", "5:48 – 6:48 PM")
    idx = 3
    for name, title, club, photo in council_except("PHF. DRR. Dr. Karishma Awari", DGS_NAME):
        slide_member(prs, name, title, club, photo, "5:48 – 6:48 PM", idx)
        idx += 1

    slide_moment(prs, "Oath Ceremony of Council", "The District Council takes its pledge", "6:48 – 6:50 PM", "Oath")

    slide_speaker(prs, "Introduction of Esteemed Guest", "DG Rtn. Nitin Dhamale", "District Governor", "6:50 – 6:53 PM")
    slide_speaker(prs, "Address by Esteemed Guest", "DG Rtn. Nitin Dhamale", "District Governor", "6:53 – 7:08 PM")
    slide_speaker(prs, "Felicitation of Esteemed Guest", "DG Rtn. Nitin Dhamale", "District Governor", "7:08 – 7:10 PM")
    slide_speaker(prs, "Address & Felicitation of Guest", "DRCC Rtn. Balvir Chawla", "District Rotaract Committee Chair", "7:10 – 7:15 PM")

    slide_section(prs, "REIGN Vision", "7:15 – 8:00 PM")
    slide_reign(prs)

    slide_moment(prs, "Launch of District Website", "rotaractdistrict3131.org", "8:00 – 8:02 PM", "Launch")
    slide_moment(prs, "Launch of District Directory", "The official RIY 2026-27 directory", "8:02 – 8:07 PM", "Launch")
    slide_moment(prs, "Special Announcements", "", "8:07 – 8:12 PM", "Announcements")
    slide_moment(prs, "Sponsor Session", "With gratitude to our partners", "8:12 – 8:19 PM", "Sponsors")
    slide_speaker(prs, "Address & Felicitation of Guest", "Dist. YSD Rtn. Santosh Pardeshi", "District Youth Service Director", "8:19 – 8:24 PM")

    slide_speaker(prs, "Address & Felicitation", "Galaxy of DRRs", "Past & Present District Representatives", "8:24 – 8:39 PM")
    slide_speaker(prs, "Introduction of Chief Guest", "Past RID Dr. Mahesh Kotbagi", "Past Rotary International Director", "8:39 – 8:41 PM")
    slide_speaker(prs, "Address by Chief Guest", "Past RID Dr. Mahesh Kotbagi", "Past Rotary International Director", "8:41 – 9:01 PM")
    slide_speaker(prs, "Felicitation of Chief Guest", "Past RID Dr. Mahesh Kotbagi", "Past Rotary International Director", "9:01 – 9:03 PM")
    slide_speaker(prs, "Closing Address", "Host Committee", "Opulence 2026", "9:03 – 9:05 PM")

    slide_moment(prs, "Secretarial Announcements", "", "9:05 – 9:07 PM", "Announcements")
    slide_moment(prs, "Open Forum", "", "9:07 – 9:09 PM", "Forum")
    slide_moment(prs, "Vote of Thanks", "", "9:09 – 9:11 PM", "Gratitude")

    slide_section(prs, "Opulence Installation\nAdjourned", "9:11 – 9:13 PM")
    slide_moment(prs, "Dinner", "Thank you for being part of Opulence", "9:13 PM onwards", "Fellowship")

    out = _resolve_out(OUT_FILE)
    prs.save(str(out))
    print(f"Created {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
